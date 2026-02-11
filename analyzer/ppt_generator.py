"""
PowerPoint Generator Service for Paper Analyzer

This service generates PowerPoint presentations from analyzed research papers,
following professional academic presentation standards with proper formatting and structure.
"""

import os
import logging
from typing import Dict, Any, Optional, Tuple
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsmap
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
import re

logger = logging.getLogger(__name__)


class PowerPointGenerator:
    """
    Generate professional academic PowerPoint presentations from paper analysis data.

    Creates a comprehensive presentation following the academic template:
    - Cover slide (title slide)
    - Introduction (1 slide)
    - Literature Review / Related Work (1 slide)
    - Problem Statement / Motivation (1 slide)
    - Main Idea & Contributions (1 slide)
    - Methodology Overview (1 slide)
    - Technical Details (1 slide)
    - Experimental Setup (1 slide)
    - Results - Quantitative (1 slide)
    - Results - Qualitative Analysis (1 slide)
    - Discussion (1 slide)
    - Conclusion & Future Work (1 slide)
    - Thank You / Q&A (1 slide)

    Total: 12-14 slides
    """

    # Professional academic color scheme (blue & teal gradient)
    PRIMARY_COLOR = RGBColor(0, 71, 160)       # Academic blue
    SECONDARY_COLOR = RGBColor(0, 122, 204)   # Light blue
    ACCENT_COLOR = RGBColor(16, 185, 129)     # Teal accent
    TEXT_COLOR = RGBColor(33, 33, 33)         # Dark gray
    LIGHT_GRAY = RGBColor(248, 250, 252)     # Very light gray
    WHITE = RGBColor(255, 255, 255)
    DARK_BLUE = RGBColor(0, 51, 102)          # Dark blue for headers

    def __init__(self, analysis_data: Dict[str, Any], metadata: Dict[str, Any]):
        """
        Initialize the PowerPoint generator.

        Args:
            analysis_data: Dictionary containing paper analysis results
            metadata: Dictionary containing paper metadata (title, authors, venue, etc.)
        """
        self.analysis_data = analysis_data
        self.metadata = metadata
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def generate(self) -> Presentation:
        """
        Generate the complete PowerPoint presentation following the academic template.

        Returns:
            Presentation: The complete PowerPoint presentation object
        """
        logger.info("Starting PowerPoint generation...")
        logger.info(f"Analysis data keys: {list(self.analysis_data.keys())}")
        logger.info(f"Metadata keys: {list(self.metadata.keys())}")

        # Log content availability
        for key, value in self.analysis_data.items():
            if isinstance(value, str):
                logger.info(f"  {key}: {len(value)} chars - {'EMPTY' if not value.strip() else 'has content'}")
            else:
                logger.info(f"  {key}: {type(value)}")

        # Slide 1: Cover slide
        logger.info("Adding cover slide...")
        self._add_cover_slide()

        # Slide 2: Introduction
        logger.info("Adding introduction slide...")
        self._add_introduction_slide()

        # Slide 3: Literature Review / Related Work
        logger.info("Adding literature review slide...")
        self._add_literature_review_slide()

        # Slide 4: Problem Statement / Motivation
        logger.info("Adding motivation slide...")
        self._add_motivation_slide()

        # Slide 5: Main Idea & Contributions
        logger.info("Adding main idea & contributions slide...")
        self._add_main_idea_slide()

        # Slide 6: Methodology Overview
        logger.info("Adding methodology overview slide...")
        self._add_methodology_overview_slide()

        # Slide 7: Technical Details
        logger.info("Adding technical details slide...")
        self._add_technical_details_slide()

        # Slide 8: Experimental Setup
        logger.info("Adding experimental setup slide...")
        self._add_experimental_setup_slide()

        # Slide 9: Results - Quantitative
        logger.info("Adding quantitative results slide...")
        self._add_quantitative_results_slide()

        # Slide 10: Results - Qualitative Analysis
        logger.info("Adding qualitative analysis slide...")
        self._add_qualitative_analysis_slide()

        # Slide 11: Discussion
        logger.info("Adding discussion slide...")
        self._add_discussion_slide()

        # Slide 12: Conclusion & Future Work
        logger.info("Adding conclusion slide...")
        self._add_conclusion_slide()

        # Slide 13: Thank You / Q&A
        logger.info("Adding thank you slide...")
        self._add_thank_you_slide()

        logger.info(f"PowerPoint generation complete. Total slides: {len(self.prs.slides)}")
        return self.prs

    def _add_cover_slide(self):
        """Create the cover slide with paper metadata."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Add background color strip at top
        left = Inches(0)
        top = Inches(0)
        width = Inches(10)
        height = Inches(1.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY_COLOR
        shape.line.fill.background()

        # Add paper title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = self.metadata.get('title', 'Paper Title')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Authors
        authors = self.metadata.get('authors', 'Unknown Authors')
        if len(authors) > 80:
            authors = authors[:77] + '...'
        authors_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(0.4))
        authors_frame = authors_box.text_frame
        authors_frame.text = f"Authors: {authors}"
        authors_para = authors_frame.paragraphs[0]
        authors_para.font.size = Pt(16)
        authors_para.font.color.rgb = self.TEXT_COLOR
        authors_para.alignment = PP_ALIGN.CENTER

        # Venue and Year
        venue = self.metadata.get('venue', 'Unknown Venue')
        year = self.metadata.get('year', 'Unknown Year')
        venue_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(0.3))
        venue_frame = venue_box.text_frame
        venue_frame.text = f"{venue} {year}"
        venue_para = venue_frame.paragraphs[0]
        venue_para.font.size = Pt(14)
        venue_para.font.color.rgb = self.SECONDARY_COLOR
        venue_para.alignment = PP_ALIGN.CENTER

        # Paper URL
        if self.metadata.get('paper_url'):
            url_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.3))
            url_frame = url_box.text_frame
            url_frame.text = self.metadata['paper_url']
            url_para = url_frame.paragraphs[0]
            url_para.font.size = Pt(10)
            url_para.font.color.rgb = RGBColor(102, 102, 102)
            url_para.alignment = PP_ALIGN.CENTER

        # Divider line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.8), Inches(6), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.SECONDARY_COLOR
        line.line.fill.background()

        # Student Information
        student_name = self.metadata.get('student_name', 'Your Name')
        student_id = self.metadata.get('student_id', 'Student ID')
        student_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
        student_frame = student_box.text_frame
        student_frame.text = f"Presented by: {student_name} ({student_id})"
        student_para = student_frame.paragraphs[0]
        student_para.font.size = Pt(18)
        student_para.font.color.rgb = self.TEXT_COLOR
        student_para.alignment = PP_ALIGN.CENTER

    def _add_title_slide(self, title: str, subtitle: str = ""):
        """
        Create a standard title slide with header and content area.

        Args:
            title: Slide title
            subtitle: Optional subtitle

        Returns:
            Tuple: (slide, content_text_frame)
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Add header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY_COLOR
        header.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.4))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(18)
            subtitle_para.font.color.rgb = self.TEXT_COLOR

        # Add content text box
        content_top = Inches(1.7) if subtitle else Inches(1.3)
        content_box = slide.shapes.add_textbox(Inches(0.5), content_top, Inches(9), Inches(5.3))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        return slide, content_frame

    def _add_introduction_slide(self):
        """Add Introduction slide."""
        introduction = self.analysis_data.get('introduction', self.analysis_data.get('abstract', ''))
        logger.info(f"Introduction content length: {len(introduction) if introduction else 0}")

        if introduction and introduction.strip():
            _, content_frame = self._add_title_slide("Introduction")
            self._add_content_to_frame(content_frame, introduction, max_slides=1)
        else:
            logger.warning("Introduction is empty, skipping slide")

    def _add_literature_review_slide(self):
        """Add Literature Review / Related Work slide."""
        literature_review = self.analysis_data.get('literature_review', '')
        if not literature_review:
            # Fallback to introduction content for related work
            introduction = self.analysis_data.get('introduction', '')
            if 'Related Work' in introduction or '### Related Work' in introduction:
                # Extract related work section
                parts = introduction.split('### Related Work')
                if len(parts) > 1:
                    literature_review = "### Related Work\n" + parts[1].split('###')[0]

        logger.info(f"Literature review content length: {len(literature_review) if literature_review else 0}")

        if literature_review and literature_review.strip():
            _, content_frame = self._add_title_slide("Literature Review")
            self._add_content_to_frame(content_frame, literature_review, max_slides=1)
        else:
            logger.warning("Literature review is empty, skipping slide")

    def _add_motivation_slide(self):
        """Add Problem Statement / Motivation slide."""
        motivation = self.analysis_data.get('motivation', '')
        logger.info(f"Motivation content length: {len(motivation) if motivation else 0}")

        if motivation and motivation.strip():
            _, content_frame = self._add_title_slide("Problem Statement / Motivation")
            self._add_content_to_frame(content_frame, motivation, max_slides=1)
        else:
            logger.warning("Motivation is empty, skipping slide")

    def _add_main_idea_slide(self):
        """Add Main Idea & Contributions slide."""
        contribution = self.analysis_data.get('contribution', '')
        logger.info(f"Main idea content length: {len(contribution) if contribution else 0}")

        if contribution and contribution.strip():
            _, content_frame = self._add_title_slide("Main Idea & Contributions")
            self._add_content_to_frame(content_frame, contribution, max_slides=1)
        else:
            logger.warning("Main idea content is empty, skipping slide")

    def _add_methodology_overview_slide(self):
        """Add Methodology Overview slide."""
        methodology = self.analysis_data.get('how_does_paper_do', '')
        logger.info(f"Methodology overview content length: {len(methodology) if methodology else 0}")

        if methodology and methodology.strip():
            # Extract the overview section (up to first ###)
            sections = methodology.split('###')
            overview = sections[0] if sections else methodology

            _, content_frame = self._add_title_slide("Methodology Overview")
            self._add_content_to_frame(content_frame, overview, max_slides=1)
        else:
            logger.warning("Methodology overview is empty, skipping slide")

    def _add_technical_details_slide(self):
        """Add Technical Details slide."""
        methodology = self.analysis_data.get('how_does_paper_do', '')
        logger.info(f"Technical details content length: {len(methodology) if methodology else 0}")

        if methodology and methodology.strip():
            # Extract sections after the first ### (Technical Details, Key Components, etc.)
            sections = methodology.split('###')
            if len(sections) > 1:
                technical_content = '###'.join(sections[1:])  # Skip the overview section
            else:
                technical_content = methodology

            _, content_frame = self._add_title_slide("Technical Details")
            self._add_content_to_frame(content_frame, technical_content, max_slides=1)
        else:
            logger.warning("Technical details is empty, skipping slide")

    def _add_experimental_setup_slide(self):
        """Add Experimental Setup slide."""
        experiments = self.analysis_data.get('what_does_paper_do', '')
        logger.info(f"Experimental setup content length: {len(experiments) if experiments else 0}")

        if experiments and experiments.strip():
            # Extract the setup section (Datasets, Metrics, Baselines)
            lines = experiments.split('\n')
            setup_lines = []
            in_setup = True
            for line in lines:
                # Stop at results section
                if '### Key Results' in line or '### Key Quantitative Results' in line or '### Main Findings' in line:
                    break
                if in_setup:
                    setup_lines.append(line)

            setup_content = '\n'.join(setup_lines) if setup_lines else experiments

            _, content_frame = self._add_title_slide("Experimental Setup")
            self._add_content_to_frame(content_frame, setup_content, max_slides=1)
        else:
            logger.warning("Experimental setup is empty, skipping slide")

    def _add_quantitative_results_slide(self):
        """Add Results - Quantitative slide."""
        experiments = self.analysis_data.get('what_does_paper_do', '')
        logger.info(f"Quantitative results content length: {len(experiments) if experiments else 0}")

        if experiments and experiments.strip():
            # Extract the results section
            lines = experiments.split('\n')
            results_lines = []
            in_results = False
            for line in lines:
                if '### Key Results' in line or '### Key Quantitative Results' in line:
                    in_results = True
                if in_results:
                    results_lines.append(line)

            results_content = '\n'.join(results_lines) if results_lines else experiments

            _, content_frame = self._add_title_slide("Results - Quantitative")
            self._add_content_to_frame(content_frame, results_content, max_slides=1)
        else:
            logger.warning("Quantitative results is empty, skipping slide")

    def _add_qualitative_analysis_slide(self):
        """Add Results - Qualitative Analysis slide."""
        results_analysis = self.analysis_data.get('results_analysis', '')
        experiments = self.analysis_data.get('what_does_paper_do', '')

        logger.info(f"Qualitative analysis content length: {len(results_analysis) if results_analysis else 0}")

        content = results_analysis if results_analysis else ""
        # Also add main findings from experiments if available
        if '### Main Findings' in experiments:
            parts = experiments.split('### Main Findings')
            if len(parts) > 1:
                findings = parts[1].split('###')[0]
                if content:
                    content += "\n\n### Main Findings\n" + findings
                else:
                    content = "### Main Findings\n" + findings

        if content and content.strip():
            _, content_frame = self._add_title_slide("Results - Qualitative Analysis")
            self._add_content_to_frame(content_frame, content, max_slides=1)
        else:
            logger.warning("Qualitative analysis is empty, skipping slide")

    def _add_discussion_slide(self):
        """Add Discussion slide."""
        discussion = self.analysis_data.get('discussion', '')
        logger.info(f"Discussion content length: {len(discussion) if discussion else 0}")

        if discussion and discussion.strip():
            _, content_frame = self._add_title_slide("Discussion")
            self._add_content_to_frame(content_frame, discussion, max_slides=1)
        else:
            logger.warning("Discussion is empty, skipping slide")

    def _add_conclusion_slide(self):
        """Add Conclusion & Future Work slide."""
        conclusion = self.analysis_data.get('conclusion', '')
        logger.info(f"Conclusion content length: {len(conclusion) if conclusion else 0}")

        if conclusion and conclusion.strip():
            _, content_frame = self._add_title_slide("Conclusion & Future Work")
            self._add_content_to_frame(content_frame, conclusion, max_slides=1)
        else:
            logger.warning("Conclusion is empty, skipping slide")

    def _add_thank_you_slide(self):
        """Add Thank You / Q&A slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Add header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY_COLOR
        header.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Thank You!"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Questions & Answers"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = self.TEXT_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Add student info
        student_name = self.metadata.get('student_name', 'Your Name')
        student_id = self.metadata.get('student_id', 'Student ID')
        contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
        contact_frame = contact_box.text_frame
        contact_frame.text = f"Presented by: {student_name} ({student_id})"
        contact_para = contact_frame.paragraphs[0]
        contact_para.font.size = Pt(18)
        contact_para.font.color.rgb = self.SECONDARY_COLOR
        contact_para.alignment = PP_ALIGN.CENTER

        # Add paper reference
        if self.metadata.get('title'):
            ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
            ref_frame = ref_box.text_frame
            ref_frame.word_wrap = True
            ref_frame.text = f"Paper: {self.metadata['title']}"
            if self.metadata.get('authors'):
                ref_frame.text += f"\nAuthors: {self.metadata['authors']}"
            if self.metadata.get('venue'):
                ref_frame.text += f"\nVenue: {self.metadata['venue']}"
            ref_para = ref_frame.paragraphs[0]
            ref_para.font.size = Pt(14)
            ref_para.font.color.rgb = self.TEXT_COLOR
            ref_para.alignment = PP_ALIGN.CENTER

    def _add_content_to_frame(self, text_frame, content: str, max_slides: int = 2):
        """
        Add markdown-formatted content to a text frame, handling overflow across slides.

        Args:
            text_frame: The text frame to add content to
            content: Markdown-formatted content string
            max_slides: Maximum number of slides to create for overflow
        """
        self._add_formatted_text(text_frame, content)

        # Check if content overflows and create additional slides if needed
        current_slide = 0
        while text_frame.text.endswith('...') and current_slide < max_slides - 1:
            # Create a continuation slide
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

            # Add header
            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
            header.fill.solid()
            header.fill.fore_color.rgb = self.PRIMARY_COLOR
            header.line.fill.background()

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = "(Continued)"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)

            # Add content text box
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            self._add_formatted_text(text_frame, content)
            current_slide += 1

    def _add_formatted_text(self, text_frame, content: str):
        """
        Parse and add markdown-formatted content to a text frame.
        Optimized for presentation readability with short, concise bullet points.

        Args:
            text_frame: PowerPoint text frame
            content: Markdown-formatted string
        """
        if not content:
            logger.warning("_add_formatted_text called with empty content")
            return

        logger.debug(f"Adding formatted text, content length: {len(content)}")
        lines = content.split('\n')
        current_paragraph = None
        paragraphs_added = 0
        max_bullets_per_slide = 6  # Limit bullets for better readability
        bullet_count = 0

        for line in lines:
            stripped = line.strip()

            # Empty line - add spacing
            if not stripped:
                if current_paragraph is not None:
                    p = text_frame.add_paragraph()
                    p.text = ""
                    p.space_after = Pt(8)
                    current_paragraph = p
                    paragraphs_added += 1
                continue

            # Header (###)
            if stripped.startswith('###'):
                header_text = stripped[3:].strip()
                p = text_frame.add_paragraph()
                p.text = header_text
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = self.PRIMARY_COLOR
                p.space_before = Pt(16)
                p.space_after = Pt(10)
                current_paragraph = p
                paragraphs_added += 1
                bullet_count = 0  # Reset bullet count after header
                logger.debug(f"Added header: {header_text[:50]}...")
                continue

            # Bullet point
            if stripped.startswith('- ') or stripped.startswith('* '):
                # Stop if we have too many bullets for readability
                if bullet_count >= max_bullets_per_slide:
                    logger.debug(f"Reached max bullets ({max_bullets_per_slide}), stopping")
                    break

                bullet_text = stripped[2:].strip()
                # Truncate long bullets for presentation readability
                if len(bullet_text) > 100:
                    bullet_text = bullet_text[:97] + "..."
                # Handle inline formatting
                bullet_text = self._parse_inline_formatting(bullet_text)

                p = text_frame.add_paragraph()
                p.text = bullet_text
                p.level = 0
                p.font.size = Pt(18)  # Larger font for better readability
                p.font.color.rgb = self.TEXT_COLOR
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.line_spacing = 1.3  # Add line spacing for readability
                current_paragraph = p
                paragraphs_added += 1
                bullet_count += 1
                logger.debug(f"Added bullet: {bullet_text[:50]}...")
                continue

            # Numbered list
            numbered_match = re.match(r'^(\d+)\.\s+(.*)$', stripped)
            if numbered_match:
                # Stop if we have too many bullets for readability
                if bullet_count >= max_bullets_per_slide:
                    logger.debug(f"Reached max bullets ({max_bullets_per_slide}), stopping")
                    break

                num_text = numbered_match.group(2)
                # Truncate long numbered items for presentation readability
                if len(num_text) > 100:
                    num_text = num_text[:97] + "..."
                # Handle inline formatting
                num_text = self._parse_inline_formatting(num_text)

                p = text_frame.add_paragraph()
                p.text = num_text
                p.level = 0
                p.font.size = Pt(18)  # Larger font for better readability
                p.font.color.rgb = self.TEXT_COLOR
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.line_spacing = 1.3  # Add line spacing for readability
                current_paragraph = p
                paragraphs_added += 1
                bullet_count += 1
                logger.debug(f"Added numbered item: {num_text[:50]}...")
                continue

            # Regular paragraph - truncate for presentation
            para_text = self._parse_inline_formatting(stripped)
            # Truncate long paragraphs for presentation readability
            if len(para_text) > 150:
                para_text = para_text[:147] + "..."

            # If it's a continuation of previous content, append to last paragraph
            if current_paragraph and len(stripped) > 0 and not stripped[0].isupper():
                current_paragraph.text += " " + para_text
            else:
                p = text_frame.add_paragraph()
                p.text = para_text
                p.font.size = Pt(18)  # Larger font for better readability
                p.font.color.rgb = self.TEXT_COLOR
                p.space_before = Pt(8)
                p.space_after = Pt(8)
                p.line_spacing = 1.4  # Add line spacing for readability
                p.alignment = PP_ALIGN.LEFT  # Left align is better for presentations
                current_paragraph = p
                paragraphs_added += 1
                bullet_count += 1
                logger.debug(f"Added paragraph: {para_text[:50]}...")

        logger.info(f"Added {paragraphs_added} paragraphs to text frame")

    def _parse_inline_formatting(self, text: str) -> str:
        """
        Parse inline markdown formatting (bold, italic, code).

        Args:
            text: Text with markdown formatting

        Returns:
            str: Text with PowerPoint formatting applied (as plain text for now,
                 as python-pptx has limited inline formatting support)
        """
        # For simplicity, we return the text stripped of markdown markers
        # In a production system, you would use run objects for proper formatting
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)  # Bold
        text = re.sub(r'\*(.*?)\*', r'\1', text)      # Italic
        text = re.sub(r'`(.*?)`', r'\1', text)        # Code
        return text

    def save(self, filepath: str):
        """
        Save the presentation to a file.

        Args:
            filepath: Path where to save the .pptx file
        """
        self.prs.save(filepath)

    def save_to_bytes(self) -> BytesIO:
        """
        Save the presentation to a BytesIO buffer.

        Returns:
            BytesIO: Buffer containing the .pptx file
        """
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer


def generate_powerpoint(analysis_data: Dict[str, Any], metadata: Dict[str, Any]) -> BytesIO:
    """
    Convenience function to generate a PowerPoint presentation.

    Args:
        analysis_data: Dictionary containing paper analysis results
        metadata: Dictionary containing paper metadata (title, authors, venue, year, url,
                       student_name, student_id)

    Returns:
        BytesIO: Buffer containing the .pptx file

    Raises:
        ValueError: If analysis_data is empty or invalid
        Exception: If PowerPoint generation fails
    """
    logger.info("=" * 80)
    logger.info("Starting PowerPoint generation process")
    logger.info(f"Analysis data keys: {list(analysis_data.keys())}")
    logger.info(f"Metadata keys: {list(metadata.keys())}")

    # Validate analysis data
    if not analysis_data:
        raise ValueError("Analysis data is empty")

    # Check if we have at least some content
    content_fields = ['introduction', 'literature_review', 'motivation', 'contribution',
                     'how_does_paper_do', 'what_does_paper_do', 'results_analysis',
                     'discussion', 'conclusion', 'future_work']
    has_content = any(analysis_data.get(field) for field in content_fields)

    if not has_content:
        logger.error("No content found in analysis data!")
        for field in content_fields:
            value = analysis_data.get(field)
            logger.error(f"  {field}: {repr(value)[:100]}")
        raise ValueError("Analysis data contains no content for slides")

    # Validate metadata
    if not metadata.get('title'):
        logger.warning("No title in metadata, using default")
        metadata['title'] = 'Untitled Presentation'

    if not metadata.get('student_name'):
        metadata['student_name'] = 'Your Name'

    if not metadata.get('student_id'):
        metadata['student_id'] = 'Student ID'

    try:
        generator = PowerPointGenerator(analysis_data, metadata)
        presentation = generator.generate()
        buffer = generator.save_to_bytes()

        # Verify the buffer has content
        buffer_size = len(buffer.getvalue())
        logger.info(f"PowerPoint generated successfully. Size: {buffer_size} bytes")

        if buffer_size < 1000:
            logger.warning(f"Generated PPT is suspiciously small ({buffer_size} bytes)")

        return buffer

    except Exception as e:
        logger.error(f"Error generating PowerPoint: {str(e)}", exc_info=True)
        raise Exception(f"PowerPoint generation failed: {str(e)}")


def extract_metadata_from_paper(paper_text: str, filename: str = "") -> Dict[str, str]:
    """
    Extract paper metadata from the paper text.

    Args:
        paper_text: Full text of the research paper
        filename: Original filename (used as fallback for title)

    Returns:
        dict: Extracted metadata (title, authors, venue, year, url)
    """
    import re

    metadata = {
        'title': filename.replace('.pdf', '') if filename else 'Unknown Title',
        'authors': 'Unknown Authors',
        'venue': 'Unknown Venue',
        'year': 'Unknown Year',
        'paper_url': '',
    }

    # Try to extract title (usually first significant text, often all caps or title case)
    lines = paper_text.split('\n')
    for i, line in enumerate(lines[:20]):  # Check first 20 lines
        stripped = line.strip()
        # Skip page numbers, headers, etc.
        if stripped and len(stripped) > 10 and len(stripped) < 200:
            # Likely a title
            if not re.match(r'^\d+$', stripped) and not stripped.startswith('http'):
                # Check if it looks like a title (contains letters, maybe some numbers)
                if any(c.isalpha() for c in stripped):
                    # First significant line is likely the title
                    if not metadata['title'] or metadata['title'] == filename.replace('.pdf', ''):
                        metadata['title'] = stripped
                    break

    # Try to extract authors (usually after title)
    in_authors = False
    author_lines = []
    for i, line in enumerate(lines[5:30]):  # Check lines 5-30
        stripped = line.strip()
        # Look for author patterns (names, universities, emails)
        if in_authors:
            if not stripped or stripped.startswith('Abstract') or len(author_lines) > 5:
                break
            if any(marker in stripped.lower() for marker in ['university', 'institute', 'college', 'lab', '@']):
                author_lines.append(stripped)
            # Check if it looks like a name (has common name patterns)
            elif any(c.isalpha() for c in stripped) and len(stripped) < 100:
                author_lines.append(stripped)
        # Start collecting after title-like text
        elif metadata['title'] in line and 'author' in lines[i+5:i+10]:
            in_authors = True

    if author_lines:
        # Join author lines and clean up
        authors_text = ' '.join(author_lines)
        # Remove common non-author text
        authors_text = re.sub(r'\bUniversity\b.*', '', authors_text, flags=re.IGNORECASE)
        authors_text = re.sub(r'\bInstitute\b.*', '', authors_text, flags=re.IGNORECASE)
        metadata['authors'] = authors_text.strip()[:200]  # Limit length

    # Try to extract venue and year from text
    text_lower = paper_text.lower()

    # Look for conference/journal names
    venues = [
        'CVPR', 'ICCV', 'ECCV', 'NeurIPS', 'ICML', 'ICLR', 'AAAI',
        'IJCAI', 'ACL', 'EMNLP', 'NAACL', 'SIGGRAPH', 'SIGMOD',
        'VLDB', 'ICDE', 'WWW', 'KDD', 'RecSys', 'CIKM',
        'IEEE Transactions on', 'ACM Transactions on',
        'Journal of', 'Proceedings of'
    ]

    for venue in venues:
        if venue.lower() in text_lower:
            # Extract the full venue mention
            pattern = re.compile(re.escape(venue) + r'[^.\n]*', re.IGNORECASE)
            matches = pattern.findall(paper_text)
            if matches:
                metadata['venue'] = matches[0].strip()[:100]
                break

    # Extract year (4 digits, likely 2015-2025)
    year_match = re.search(r'\b(201[5-9]|202[0-5])\b', paper_text[:5000])
    if year_match:
        metadata['year'] = year_match.group(1)

    # Try to extract arXiv URL or other URLs
    url_match = re.search(r'(https?://arxiv\.org/abs/\d+\.\d+)', paper_text[:3000])
    if url_match:
        metadata['paper_url'] = url_match.group(1)
    else:
        # Look for any other paper URL
        url_match = re.search(r'(https?://[^\s]+(?:\.pdf|/paper/))', paper_text[:3000])
        if url_match:
            metadata['paper_url'] = url_match.group(1)

    return metadata
